from pathlib import Path
from typing import Dict, Tuple

from pptx import Presentation


def extract_pptx_text(file_path: Path) -> Tuple[str, Dict[str, str]]:
    presentation = Presentation(str(file_path))
    collected = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    collected.append(text)
    combined = "\n".join(collected).strip()
    return combined, {"slides": str(len(presentation.slides))}
