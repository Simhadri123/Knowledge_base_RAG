from pathlib import Path
import sys

from fastapi.testclient import TestClient
from PIL import Image, ImageDraw
import fitz
from docx import Document
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parents[1]
if str(BASE_DIR) not in sys.path:
    sys.path.insert(0, str(BASE_DIR))

from app.main import app


def create_sample_vtt(path: Path) -> None:
    content = """WEBVTT

00:00:00.000 --> 00:00:02.000
Hello from the transcript.

00:00:02.000 --> 00:00:04.000
This is another line.
"""
    path.write_text(content, encoding="utf-8")


def create_sample_docx(path: Path) -> None:
    doc = Document()
    doc.add_paragraph("Sample DOCX text for extraction.")
    doc.save(str(path))


def create_sample_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Sample PPTX"
    slide.placeholders[1].text = "Slide body text for extraction."
    presentation.save(str(path))


def create_sample_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Sample PDF text for extraction.")
    document.save(str(path))


def create_sample_image(path: Path) -> None:
    image = Image.new("RGB", (400, 120), color="white")
    drawer = ImageDraw.Draw(image)
    drawer.text((10, 40), "Sample OCR text", fill="black")
    image.save(str(path))


def run_uploads(sample_dir: Path) -> None:
    client = TestClient(app)
    samples = {
        "sample.vtt": create_sample_vtt,
        "sample.docx": create_sample_docx,
        "sample.pptx": create_sample_pptx,
        "sample.pdf": create_sample_pdf,
        "sample.png": create_sample_image,
        "sample.jpg": create_sample_image,
    }

    for name, creator in samples.items():
        file_path = sample_dir / name
        creator(file_path)
        with file_path.open("rb") as handle:
            response = client.post("/upload", files={"file": (name, handle)})

        assert response.status_code == 200, response.text
        payload = response.json()
        extracted_text = payload.get("text", "")
        print(f"\n{name} extracted text:\n{extracted_text}\n")
        assert extracted_text.strip(), f"Empty extraction for {name}"


if __name__ == "__main__":
    temp_dir = Path(__file__).parent / "_samples"
    temp_dir.mkdir(parents=True, exist_ok=True)
    run_uploads(temp_dir)
    print("All Phase 1 extraction tests passed.")
